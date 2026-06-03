from __future__ import annotations

from datetime import datetime
from io import BytesIO, StringIO
import re
import zipfile

import pandas as pd
import streamlit as st

try:
    from PIL import Image
except Exception:  # pragma: no cover - optional dependency fallback
    Image = None

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
except Exception:  # pragma: no cover - optional dependency fallback
    colors = None
    A4 = None
    getSampleStyleSheet = None
    cm = None
    Paragraph = None
    SimpleDocTemplate = None
    Spacer = None
    Table = None
    TableStyle = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover - optional dependency fallback
    Presentation = None
    RGBColor = None
    MSO_SHAPE = None
    PP_ALIGN = None
    Inches = None
    Pt = None


TAIVAS_TARGET_FIELDS = [
    "timestamp",
    "country",
    "city",
    "lat",
    "lon",
    "temperature",
    "wind_speed",
    "solar_radiation",
    "precipitation",
    "humidity",
    "population",
    "solar_capacity",
    "wind_capacity",
    "geothermal_capacity",
    "hydro_capacity",
    "battery_capacity",
]

TAIVAS_RECOMMENDED_FIELDS = [
    "timestamp",
    "country",
    "city",
    "temperature",
    "wind_speed",
    "solar_radiation",
    "precipitation",
    "humidity",
    "population",
    "solar_capacity",
    "wind_capacity",
    "geothermal_capacity",
    "hydro_capacity",
    "battery_capacity",
]

TAIVAS_NUMERIC_FIELDS = [
    "lat",
    "lon",
    "temperature",
    "wind_speed",
    "solar_radiation",
    "precipitation",
    "humidity",
    "population",
    "solar_capacity",
    "wind_capacity",
    "geothermal_capacity",
    "hydro_capacity",
    "battery_capacity",
]

NON_NEGATIVE_FIELDS = [
    "wind_speed",
    "solar_radiation",
    "precipitation",
    "humidity",
    "population",
    "solar_capacity",
    "wind_capacity",
    "geothermal_capacity",
    "hydro_capacity",
    "battery_capacity",
]

COLUMN_ALIASES = {
    "time": "timestamp",
    "date": "timestamp",
    "datetime": "timestamp",
    "timestamp_utc": "timestamp",
    "country_name": "country",
    "nation": "country",
    "city_name": "city",
    "location": "city",
    "latitude": "lat",
    "longitude": "lon",
    "long": "lon",
    "temp": "temperature",
    "temp_c": "temperature",
    "temperature_c": "temperature",
    "wind": "wind_speed",
    "wind_m_s": "wind_speed",
    "wind_speed_m_s": "wind_speed",
    "solar": "solar_radiation",
    "irradiance": "solar_radiation",
    "solar_irradiance": "solar_radiation",
    "rain": "precipitation",
    "rainfall": "precipitation",
    "precip": "precipitation",
    "humid": "humidity",
    "pop": "population",
    "solar_mw": "solar_capacity",
    "wind_mw": "wind_capacity",
    "geo_capacity": "geothermal_capacity",
    "geothermal_mw": "geothermal_capacity",
    "hydro_mw": "hydro_capacity",
    "battery_mwh": "battery_capacity",
    "storage_capacity": "battery_capacity",
}

PDF_DISCLAIMER = (
    "TAIVAS is a decision-support simulation tool. It does not provide guaranteed "
    "predictions, emergency commands, or final operational decisions. Results should "
    "be reviewed together with real-time data, professional judgment, and institutional protocols."
)


def normalize_column_name(name: object) -> str:
    normalized = str(name).strip().lower()
    normalized = re.sub(r"[^a-z0-9]+", "_", normalized).strip("_")
    return COLUMN_ALIASES.get(normalized, normalized)


def normalize_taivas_dataframe(df: pd.DataFrame) -> tuple[pd.DataFrame, list[str]]:
    cleaned = df.copy()
    cleaned.columns = [normalize_column_name(col) for col in cleaned.columns]
    duplicate_columns = cleaned.columns[cleaned.columns.duplicated()].tolist()
    if duplicate_columns:
        cleaned = cleaned.loc[:, ~cleaned.columns.duplicated()]
    ordered_columns = [col for col in TAIVAS_TARGET_FIELDS if col in cleaned.columns]
    extra_columns = [col for col in cleaned.columns if col not in ordered_columns]
    cleaned = cleaned[ordered_columns + extra_columns]
    missing = [col for col in TAIVAS_TARGET_FIELDS if col not in cleaned.columns]
    warnings = []
    if duplicate_columns:
        warnings.append(f"Duplicate normalized columns were removed: {', '.join(sorted(set(duplicate_columns)))}")
    if missing:
        warnings.append(f"Missing TAIVAS fields: {', '.join(missing)}")
    return cleaned, warnings


def dataframe_to_csv_bytes(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8-sig")


def validate_taivas_csv_dataframe(df: pd.DataFrame) -> tuple[str, str, pd.DataFrame, pd.DataFrame]:
    cleaned, normalize_warnings = normalize_taivas_dataframe(df)
    issues = []
    for warning in normalize_warnings:
        issues.append({"Severity": "Warning", "Field": "columns", "Issue": warning})

    if cleaned.empty:
        issues.append({"Severity": "Error", "Field": "file", "Issue": "The uploaded CSV has no rows."})

    for field in TAIVAS_RECOMMENDED_FIELDS:
        if field not in cleaned.columns:
            issues.append({"Severity": "Warning", "Field": field, "Issue": "Recommended field is missing."})

    for field in cleaned.columns:
        empty_count = int(cleaned[field].isna().sum() + (cleaned[field].astype(str).str.strip() == "").sum())
        if empty_count > 0:
            severity = "Warning" if field not in ("country", "city") else "Error"
            issues.append({"Severity": severity, "Field": field, "Issue": f"{empty_count} empty value(s) detected."})

    if "timestamp" in cleaned.columns:
        parsed = pd.to_datetime(cleaned["timestamp"], errors="coerce")
        invalid_count = int(parsed.isna().sum())
        if invalid_count > 0:
            issues.append({"Severity": "Warning", "Field": "timestamp", "Issue": f"{invalid_count} timestamp value(s) could not be parsed."})
        duplicate_count = int(cleaned["timestamp"].duplicated().sum())
        if duplicate_count > 0:
            issues.append({"Severity": "Warning", "Field": "timestamp", "Issue": f"{duplicate_count} duplicate timestamp row(s) detected."})

    for field in TAIVAS_NUMERIC_FIELDS:
        if field not in cleaned.columns:
            continue
        numeric_values = pd.to_numeric(cleaned[field], errors="coerce")
        invalid_count = int(numeric_values.isna().sum() - cleaned[field].isna().sum())
        if invalid_count > 0:
            issues.append({"Severity": "Warning", "Field": field, "Issue": f"{invalid_count} non-numeric value(s) detected."})
        if field in NON_NEGATIVE_FIELDS:
            negative_count = int((numeric_values < 0).sum())
            if negative_count > 0:
                issues.append({"Severity": "Warning", "Field": field, "Issue": f"{negative_count} negative value(s) detected."})

    ranges = {
        "temperature": (-60, 70),
        "wind_speed": (0, 80),
        "solar_radiation": (0, 1400),
        "precipitation": (0, 1000),
        "humidity": (0, 100),
        "population": (1, 100000000),
        "solar_capacity": (0, 100000),
        "wind_capacity": (0, 100000),
        "geothermal_capacity": (0, 100000),
        "hydro_capacity": (0, 100000),
        "battery_capacity": (0, 1000000),
    }
    for field, (low, high) in ranges.items():
        if field not in cleaned.columns:
            continue
        numeric_values = pd.to_numeric(cleaned[field], errors="coerce")
        suspicious_count = int(((numeric_values < low) | (numeric_values > high)).sum())
        if suspicious_count > 0:
            issues.append({"Severity": "Warning", "Field": field, "Issue": f"{suspicious_count} suspicious value(s) outside expected range {low} to {high}."})

    for field in ("country", "city"):
        if field in cleaned.columns and cleaned[field].dropna().astype(str).str.strip().nunique() > 1:
            issues.append({"Severity": "Warning", "Field": field, "Issue": f"Multiple {field} values detected. Verify the file is intended for mixed locations."})

    issues_df = pd.DataFrame(issues, columns=["Severity", "Field", "Issue"])
    has_error = not issues_df.empty and (issues_df["Severity"] == "Error").any()
    has_warning = not issues_df.empty and (issues_df["Severity"] == "Warning").any()
    if has_error:
        status = "Error"
        message = "This file is not ready. Please fix the listed issues."
    elif has_warning:
        status = "Warning"
        message = "This file can be used, but some fields need review."
    else:
        status = "Pass"
        message = "This file appears ready for TAIVAS."
    return status, message, issues_df, cleaned


def build_pdf_report(context: dict, results: dict, inputs: dict, notes: str = "") -> bytes:
    if SimpleDocTemplate is None:
        raise RuntimeError("PDF support requires reportlab. Install reportlab or update requirements.txt.")

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=1.6 * cm, leftMargin=1.6 * cm, topMargin=1.5 * cm, bottomMargin=1.5 * cm)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("TAIVAS Energy Resilience Simulation Report", styles["Title"]),
        Spacer(1, 0.25 * cm),
        Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", styles["Normal"]),
        Spacer(1, 0.35 * cm),
    ]

    summary_rows = [
        ["Country", context.get("country", "-")],
        ["City", context.get("city", "-")],
        ["Scenario", str(context.get("scenario", "-")).replace("_", " ").title()],
        ["Risk Level", context.get("risk_level", results.get("risk_tier", "-"))],
    ]
    capacity_rows = [
        ["Solar Capacity", f"{inputs.get('solar_capacity', '-')} MW"],
        ["Wind Capacity", f"{inputs.get('wind_capacity', '-')} MW"],
        ["Geothermal Capacity", f"{inputs.get('geothermal_capacity', '-')} MW"],
        ["Hydro Capacity", f"{inputs.get('hydro_capacity', '-')} MW"],
        ["Battery Capacity", f"{inputs.get('battery_capacity', '-')} MWh"],
    ]
    output_rows = [
        ["Demand", f"{results.get('demand', '-')} MW"],
        ["Renewable Supply", f"{results.get('renewable_supply', '-')} MW"],
        ["Final Supply", f"{results.get('final_supply', '-')} MW"],
        ["Possible Energy Gap", f"{results.get('shortfall', '-')} MW"],
        ["Renewable Ratio", f"{results.get('renewable_ratio', '-')}%"],
        ["Battery Level", f"{results.get('battery_levels', '-')} MWh"],
        ["System Efficiency", f"{results.get('system_efficiency', '-')}%"],
        ["Grid Dependency", f"{results.get('grid_dependency', '-')}%"],
    ]

    for title, rows in (("Scenario Summary", summary_rows), ("Input Capacity Summary", capacity_rows), ("Key Output Summary", output_rows)):
        story.append(Paragraph(title, styles["Heading2"]))
        table = Table(rows, colWidths=[5.2 * cm, 10.5 * cm])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#E8EEF7")),
            ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#111827")),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CBD5E1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("PADDING", (0, 0), (-1, -1), 6),
        ]))
        story.append(table)
        story.append(Spacer(1, 0.35 * cm))

    if notes.strip():
        story.append(Paragraph("Notes", styles["Heading2"]))
        story.append(Paragraph(notes.strip(), styles["Normal"]))
        story.append(Spacer(1, 0.25 * cm))

    story.append(Paragraph("Decision-Support Disclaimer", styles["Heading2"]))
    story.append(Paragraph(PDF_DISCLAIMER, styles["Normal"]))
    doc.build(story)
    return buffer.getvalue()


def format_metric(value: object, unit: str = "") -> str:
    if value is None or value == "":
        return "-"
    if isinstance(value, float):
        text = f"{value:.2f}"
    else:
        text = str(value)
    return f"{text} {unit}".strip()


def recommendation_from_results(results: dict) -> str:
    shortfall = float(results.get("shortfall", 0) or 0)
    battery = float(results.get("battery_levels", 0) or 0)
    grid_dependency = float(results.get("grid_dependency", 0) or 0)
    if shortfall > 0:
        return "Review critical-load prioritization, reduce non-essential demand, and verify backup supply readiness."
    if battery <= 0:
        return "Review storage assumptions and confirm available backup capacity before operational use."
    if grid_dependency >= 10:
        return "Review external power dependency and confirm backup grid support assumptions."
    return "Continue monitoring scenario assumptions and preserve reserve margin."


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=(17, 24, 39), align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.48), title, font_size=25, bold=True, color=(15, 23, 42))
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.86), Inches(12.0), Inches(0.34), subtitle, font_size=11, color=(71, 85, 105))


def add_key_value_table(slide, rows, left=0.7, top=1.35, width=12.0, height=4.6):
    usable_rows = [[str(a), str(b)] for a, b in rows if b is not None]
    if not usable_rows:
        usable_rows = [["Status", "No data available"]]
    table_shape = slide.shapes.add_table(len(usable_rows), 2, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    table.columns[0].width = Inches(width * 0.38)
    table.columns[1].width = Inches(width * 0.62)
    for row_idx, row in enumerate(usable_rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(232, 238, 247) if col_idx == 0 else RGBColor(248, 250, 252)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.bold = col_idx == 0
                    run.font.color.rgb = RGBColor(17, 24, 39)
    return table_shape


def add_metric_cards(slide, metrics):
    for index, (label, value) in enumerate(metrics[:4]):
        left = Inches(0.65 + index * 3.1)
        top = Inches(1.45)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.82), Inches(1.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
        shape.line.color.rgb = RGBColor(203, 213, 225)
        add_textbox(slide, left + Inches(0.16), top + Inches(0.18), Inches(2.45), Inches(0.3), label, font_size=10, color=(71, 85, 105))
        add_textbox(slide, left + Inches(0.16), top + Inches(0.58), Inches(2.45), Inches(0.5), value, font_size=18, bold=True, color=(15, 23, 42))


def add_bar_comparison(slide, rows):
    max_value = max([float(value or 0) for _, value, _ in rows] + [1.0])
    top = Inches(1.45)
    for index, (label, value, color) in enumerate(rows):
        value = float(value or 0)
        y = top + Inches(index * 0.72)
        add_textbox(slide, Inches(0.75), y, Inches(2.2), Inches(0.32), label, font_size=11, bold=True)
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.05), y, Inches(7.6), Inches(0.28))
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(226, 232, 240)
        background.line.color.rgb = RGBColor(226, 232, 240)
        bar_width = max(0.05, 7.6 * value / max_value)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.05), y, Inches(bar_width), Inches(0.28))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*color)
        bar.line.color.rgb = RGBColor(*color)
        add_textbox(slide, Inches(10.85), y - Inches(0.02), Inches(1.5), Inches(0.32), f"{value:.2f} MW", font_size=10)


def build_ppt_report(context: dict, results: dict, inputs: dict, notes: str = "") -> bytes:
    # Rapid briefing deck generator: stable structure and correct data mapping
    # are prioritized over advanced visual design.
    if Presentation is None:
        raise RuntimeError("PPT support requires python-pptx. Install python-pptx or update requirements.txt.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    generated = datetime.now().strftime("%Y-%m-%d %H:%M")
    country = context.get("country", "-")
    city = context.get("city", "-")
    scenario = str(context.get("scenario", "-")).replace("_", " ").title()
    risk_level = context.get("risk_level", results.get("risk_tier", "-"))

    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, Inches(0.75), Inches(1.1), Inches(11.8), Inches(0.7), "TAIVAS Rapid Briefing Deck", font_size=32, bold=True, color=(15, 23, 42))
    add_textbox(slide, Inches(0.78), Inches(1.92), Inches(11.5), Inches(0.42), f"{city}, {country} | {scenario}", font_size=18, color=(51, 65, 85))
    add_textbox(slide, Inches(0.78), Inches(2.52), Inches(11.3), Inches(0.34), f"Generated: {generated}", font_size=12, color=(100, 116, 139))
    add_textbox(slide, Inches(0.78), Inches(5.85), Inches(11.7), Inches(0.5), "Decision-support simulation only. Not a guaranteed forecast or command system.", font_size=13, bold=True, color=(71, 85, 105))

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Test Summary", "High-level context for the completed TAIVAS simulation.")
    add_key_value_table(slide, [
        ["Country", country],
        ["City", city],
        ["Scenario", scenario],
        ["Risk Level", risk_level],
        ["Generated", generated],
    ])

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Scenario Input Summary", "Installed capacity and scenario inputs used by the current TAIVAS run.")
    add_key_value_table(slide, [
        ["Solar Capacity", format_metric(inputs.get("solar_capacity"), "MW")],
        ["Wind Capacity", format_metric(inputs.get("wind_capacity"), "MW")],
        ["Geothermal Capacity", format_metric(inputs.get("geothermal_capacity"), "MW")],
        ["Hydro Capacity", format_metric(inputs.get("hydro_capacity"), "MW")],
        ["Battery Capacity", format_metric(inputs.get("battery_capacity"), "MWh")],
    ])

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Key Output Metrics", "Primary outputs from the completed simulation.")
    add_metric_cards(slide, [
        ["Demand", format_metric(results.get("demand"), "MW")],
        ["Renewable Supply", format_metric(results.get("renewable_supply"), "MW")],
        ["Final Supply", format_metric(results.get("final_supply"), "MW")],
        ["Energy Gap", format_metric(results.get("shortfall"), "MW")],
    ])
    add_key_value_table(slide, [
        ["Renewable Ratio", format_metric(results.get("renewable_ratio"), "%")],
        ["System Efficiency", format_metric(results.get("system_efficiency"), "%")],
        ["Grid Dependency", format_metric(results.get("grid_dependency"), "%")],
        ["Battery Level", format_metric(results.get("battery_levels"), "MWh")],
    ], top=3.25, height=2.5)

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Renewable Supply vs Demand", "Visual comparison of current demand and available supply.")
    add_bar_comparison(slide, [
        ["Demand", results.get("demand", 0), (56, 189, 248)],
        ["Renewable Supply", results.get("renewable_supply", 0), (34, 197, 94)],
        ["Final Supply", results.get("final_supply", 0), (96, 165, 250)],
    ])

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Possible Energy Gap Summary", "Shortfall is shown as a scenario output, not a confirmed prediction.")
    gap = float(results.get("shortfall", 0) or 0)
    demand = float(results.get("demand", 0) or 0)
    gap_pct = (gap / demand * 100) if demand else 0
    add_key_value_table(slide, [
        ["Possible Energy Gap", f"{gap:.2f} MW"],
        ["Gap Share of Demand", f"{gap_pct:.2f}%"],
        ["Interpretation", "Review non-critical load and backup supply if an energy gap appears." if gap > 0 else "No modeled energy gap under current scenario assumptions."],
    ])

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Battery Level Summary", "Storage is treated as a resilience buffer in the current simulation output.")
    add_key_value_table(slide, [
        ["Battery Level", format_metric(results.get("battery_levels"), "MWh")],
        ["Battery Capacity", format_metric(inputs.get("battery_capacity"), "MWh")],
        ["Review Note", "Verify storage assumptions and operational availability before decisions."],
    ])

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Risk Level and Recommendation", "Decision-support guidance generated from current simulation outputs.")
    add_key_value_table(slide, [
        ["Risk Level", risk_level],
        ["Recommendation", recommendation_from_results(results)],
        ["Human Review", "Human confirmation is required before operational changes."],
    ])

    scenario_rows = context.get("scenario_comparison_rows")
    if scenario_rows:
        slide = prs.slides.add_slide(blank_layout)
        add_slide_title(slide, "Scenario Comparison", "Comparison table from available TAIVAS scenario results.")
        rows = [["Scenario", "Energy Gap", "System Efficiency"]]
        for item in scenario_rows[:8]:
            rows.append([
                item.get("Scenario", "-"),
                str(item.get("Energy Gap (MW)", item.get("Energy Gap", "-"))),
                str(item.get("System Stability (%)", item.get("System Efficiency", "-"))),
            ])
        table_shape = slide.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.35), Inches(12.0), Inches(4.8))
        table = table_shape.table
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(226, 232, 240) if row_idx == 0 else RGBColor(248, 250, 252)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
                        run.font.bold = row_idx == 0

    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, "Decision-Support Disclaimer", "Scope boundary for institutional review.")
    add_textbox(slide, Inches(0.8), Inches(1.45), Inches(11.7), Inches(1.4), PDF_DISCLAIMER, font_size=17, color=(30, 41, 59))
    if notes.strip():
        add_textbox(slide, Inches(0.8), Inches(3.35), Inches(11.7), Inches(1.2), f"Notes: {notes.strip()}", font_size=13, color=(71, 85, 105))

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def safe_filename(value: object) -> str:
    text = str(value or "TAIVAS").strip()
    text = re.sub(r"[^A-Za-z0-9_-]+", "_", text).strip("_")
    return text or "TAIVAS"


def convert_image_file(uploaded_file, output_format: str, quality: int, resize_width: int | None) -> tuple[str, bytes, int, int]:
    if Image is None:
        raise RuntimeError("Image conversion requires Pillow. Install pillow or update requirements.txt.")
    source_bytes = uploaded_file.getvalue()
    with Image.open(BytesIO(source_bytes)) as image:
        image = image.convert("RGB") if output_format.upper() in ("JPEG", "JPG", "WEBP") else image.convert("RGBA")
        if resize_width and resize_width > 0 and image.width > resize_width:
            new_height = max(1, int(image.height * (resize_width / image.width)))
            image = image.resize((resize_width, new_height))
        output = BytesIO()
        save_format = "JPEG" if output_format.upper() == "JPG" else output_format.upper()
        save_kwargs = {}
        if save_format in ("JPEG", "WEBP"):
            save_kwargs["quality"] = int(quality)
            save_kwargs["optimize"] = True
        image.save(output, format=save_format, **save_kwargs)
    extension = "jpg" if output_format.upper() in ("JPG", "JPEG") else output_format.lower()
    base_name = safe_filename(".".join(uploaded_file.name.split(".")[:-1]) or uploaded_file.name)
    return f"{base_name}.{extension}", output.getvalue(), len(source_bytes), len(output.getvalue())


def zip_bytes(file_items: list[tuple[str, bytes]]) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for filename, content in file_items:
            archive.writestr(filename, content)
    return buffer.getvalue()


def render_excel_to_csv_tool():
    st.markdown("Convert an Excel sheet into a TAIVAS-compatible CSV. Missing fields are reported as warnings, not hard failures.")
    uploaded_excel = st.file_uploader("Upload Excel file", type=["xlsx", "xls"], key="toolkit_excel_upload")
    if not uploaded_excel:
        return
    try:
        excel_file = pd.ExcelFile(uploaded_excel)
    except Exception as exc:
        st.error(f"Unable to read Excel file: {exc}")
        return
    sheet_name = st.selectbox("Select sheet", excel_file.sheet_names, key="toolkit_excel_sheet")
    try:
        raw_df = pd.read_excel(uploaded_excel, sheet_name=sheet_name)
    except Exception as exc:
        st.error(f"Unable to load selected sheet: {exc}")
        return
    cleaned_df, warnings = normalize_taivas_dataframe(raw_df)
    if warnings:
        for warning in warnings:
            st.warning(warning)
    st.subheader("Preview")
    st.dataframe(cleaned_df.head(20), use_container_width=True)
    st.download_button(
        "Download TAIVAS CSV",
        dataframe_to_csv_bytes(cleaned_df),
        file_name=f"taivas_converted_{safe_filename(sheet_name)}.csv",
        mime="text/csv",
    )


def render_csv_validator_tool():
    st.markdown("Check whether a CSV file is ready for TAIVAS. The validator favors clear warnings over aggressive rejection.")
    uploaded_csv = st.file_uploader("Upload TAIVAS CSV", type=["csv"], key="toolkit_csv_validator")
    if not uploaded_csv:
        return
    try:
        df = pd.read_csv(uploaded_csv)
    except Exception as exc:
        st.error(f"Unable to read CSV file: {exc}")
        return
    status, message, issues_df, cleaned_df = validate_taivas_csv_dataframe(df)
    status_method = st.success if status == "Pass" else st.warning if status == "Warning" else st.error
    status_method(message)
    summary_cols = st.columns(3)
    summary_cols[0].metric("Status", status)
    summary_cols[1].metric("Rows", len(cleaned_df))
    summary_cols[2].metric("Issues", len(issues_df))
    if not issues_df.empty:
        st.subheader("Detected Issues")
        st.dataframe(issues_df, use_container_width=True, hide_index=True)
    st.subheader("Cleaned Preview")
    st.dataframe(cleaned_df.head(20), use_container_width=True)
    st.download_button(
        "Download Cleaned CSV",
        dataframe_to_csv_bytes(cleaned_df),
        file_name="taivas_cleaned.csv",
        mime="text/csv",
    )


def render_pdf_report_tool(current_context: dict, current_results: dict, current_inputs: dict):
    st.markdown("Generate a simple decision-support PDF or rapid PowerPoint briefing deck from the current TAIVAS simulation result.")
    st.caption("The PPT is a structured draft for review. It is not a highly designed investor deck.")
    notes = st.text_area("Optional notes", "", key="toolkit_report_notes")
    pdf_col, ppt_col = st.columns(2)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    base_filename = (
        f"TAIVAS_Report_{safe_filename(current_context.get('country'))}_"
        f"{safe_filename(current_context.get('city'))}_{safe_filename(current_context.get('scenario'))}_{timestamp}"
    )
    with pdf_col:
        if st.button("Generate PDF Report", key="toolkit_pdf_generate"):
            try:
                pdf_bytes = build_pdf_report(current_context, current_results, current_inputs, notes)
            except Exception as exc:
                st.error(f"Unable to generate PDF report: {exc}")
            else:
                st.download_button("Download PDF Report", pdf_bytes, file_name=f"{base_filename}.pdf", mime="application/pdf")
                st.success("PDF report generated.")
    with ppt_col:
        if st.button("Generate PPT Briefing Deck", key="toolkit_ppt_generate"):
            try:
                ppt_bytes = build_ppt_report(current_context, current_results, current_inputs, notes)
            except Exception as exc:
                st.error(f"Unable to generate PPT report: {exc}")
            else:
                ppt_filename = (
                    f"TAIVAS_Presentation_{safe_filename(current_context.get('country'))}_"
                    f"{safe_filename(current_context.get('city'))}_{safe_filename(current_context.get('scenario'))}_{timestamp}.pptx"
                )
                st.download_button(
                    "Download PPT Briefing Deck",
                    ppt_bytes,
                    file_name=ppt_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                st.success("PPT report generated.")
    st.caption(PDF_DISCLAIMER)


def render_image_converter_tool():
    st.markdown("Convert or compress screenshots and presentation images. HEIC may not be supported unless Pillow can read it in the deployment environment.")
    uploaded_images = st.file_uploader("Upload images", type=["png", "jpg", "jpeg", "webp", "heic"], accept_multiple_files=True, key="toolkit_image_upload")
    output_format = st.selectbox("Output format", ["PNG", "JPG", "WEBP"], key="toolkit_image_format")
    quality = st.slider("Quality for JPG/WEBP", 30, 100, 82, key="toolkit_image_quality")
    resize_enabled = st.checkbox("Resize width", value=False, key="toolkit_image_resize_enabled")
    resize_width = st.number_input("Maximum width in pixels", min_value=100, max_value=6000, value=1600, step=100, disabled=not resize_enabled, key="toolkit_image_width")
    if not uploaded_images:
        return
    converted_items = []
    for uploaded_image in uploaded_images:
        try:
            filename, content, original_size, output_size = convert_image_file(
                uploaded_image,
                output_format,
                quality,
                int(resize_width) if resize_enabled else None,
            )
        except Exception as exc:
            st.error(f"{uploaded_image.name}: {exc}")
            continue
        converted_items.append((filename, content))
        st.write(f"{uploaded_image.name} -> {filename} | {original_size / 1024:.1f} KB -> {output_size / 1024:.1f} KB")
        st.download_button(f"Download {filename}", content, file_name=filename, mime=f"image/{'jpeg' if filename.endswith('.jpg') else output_format.lower()}")
    if len(converted_items) > 1:
        st.download_button("Download All as ZIP", zip_bytes(converted_items), file_name="taivas_converted_images.zip", mime="application/zip")


def render_batch_renamer_tool():
    st.markdown("Create renamed copies of uploaded files and download them as a ZIP. Original files are not renamed on the server.")
    uploaded_files = st.file_uploader("Upload files to rename", accept_multiple_files=True, key="toolkit_rename_upload")
    prefix = st.text_input("Filename prefix", "TAIVAS_Dashboard", key="toolkit_rename_prefix")
    numbering = st.selectbox("Numbering style", ["001, 002, 003...", "01, 02, 03..."], key="toolkit_rename_numbering")
    digits = 3 if numbering.startswith("001") else 2
    if not uploaded_files:
        return
    preview = []
    renamed_items = []
    for index, uploaded_file in enumerate(uploaded_files, start=1):
        extension = uploaded_file.name.split(".")[-1] if "." in uploaded_file.name else "bin"
        new_name = f"{safe_filename(prefix)}_{index:0{digits}d}.{extension}"
        preview.append({"Original Filename": uploaded_file.name, "New Filename": new_name})
        renamed_items.append((new_name, uploaded_file.getvalue()))
    st.subheader("Preview")
    st.dataframe(pd.DataFrame(preview), use_container_width=True, hide_index=True)
    st.download_button("Download Renamed Files ZIP", zip_bytes(renamed_items), file_name=f"{safe_filename(prefix)}_renamed.zip", mime="application/zip")


def render_taivas_utility_toolkit(current_context: dict, current_results: dict, current_inputs: dict):
    st.subheader("TAIVAS Utility Toolkit")
    st.caption("Practical file utilities for TAIVAS workflows. These tools do not change simulation formulas or scenario logic.")
    tabs = st.tabs([
        "Excel to TAIVAS CSV",
        "CSV Validator",
        "PDF / PPT Report Generator",
        "Image Converter",
        "Batch Renamer",
    ])
    with tabs[0]:
        render_excel_to_csv_tool()
    with tabs[1]:
        render_csv_validator_tool()
    with tabs[2]:
        render_pdf_report_tool(current_context, current_results, current_inputs)
    with tabs[3]:
        render_image_converter_tool()
    with tabs[4]:
        render_batch_renamer_tool()
